from pptx import Presentation
from PIL import Image


class Pptx:

    def __init__(self, assunto, theme):
        self.assunto = assunto
        self.ppt = Presentation(theme)
        self.layouts = self.ppt.slide_layouts

    def add_title_slide(self):
        self.ppt.slides.add_slide(self.layouts[0]).shapes.title.text = self.assunto.title()

    def add_context(self, text, image_path):
        slide = self.ppt.slides.add_slide(self.layouts[8])

        slide.placeholders[2].text = text

        self._add_image(slide, 1, image_path)

    def add_thanks_for_watching_slide(self):
        self.ppt.slides.add_slide(self.layouts[0]).shapes.title.text = 'Obrigado por assistir'

    def save(self):
        self.ppt.save(f'Presentations/{self.assunto}.pptx')

    @classmethod
    def _add_image(cls, slide, placeholder_id, image_url):
        placeholder = slide.placeholders[placeholder_id]

        im = Image.open(image_url)
        width, height = im.size

        placeholder.height = height
        placeholder.width = width

        placeholder = placeholder.insert_picture(image_url)

        image_ratio = width / height
        placeholder_ratio = placeholder.width / placeholder.height
        ratio_difference = placeholder_ratio - image_ratio

        if ratio_difference > 0:
            difference_on_each_side = ratio_difference / 2
            placeholder.crop_left = -difference_on_each_side
            placeholder.crop_right = -difference_on_each_side

        else:
            difference_on_each_side = -ratio_difference / 2
            placeholder.crop_bottom = -difference_on_each_side
            placeholder.crop_top = -difference_on_each_side
